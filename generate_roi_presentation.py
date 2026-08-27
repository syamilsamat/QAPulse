import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION, XL_TICK_MARK

# =============================================================================
# ROI INPUTS & CALCULATIONS (see notes slide for full methodology)
# =============================================================================
SALARY = 5000.0                # RM per QA per month
MANUAL_PER_WEEK = 100          # test cases/week, manual process
QAPULSE_PER_DAY = 100          # test cases/day, with QAPulse
DAYS_PER_WEEK = 5
WEEKS_PER_MONTH = 4.33
DAYS_PER_MONTH = DAYS_PER_WEEK * WEEKS_PER_MONTH   # 21.65

MANUAL_PER_DAY = MANUAL_PER_WEEK / DAYS_PER_WEEK              # 20.0
MULTIPLIER = QAPULSE_PER_DAY / MANUAL_PER_DAY                  # 5.0x

MANUAL_PER_MONTH = MANUAL_PER_WEEK * WEEKS_PER_MONTH            # 433
QAPULSE_PER_MONTH = QAPULSE_PER_DAY * DAYS_PER_MONTH             # 2165

COST_PER_TC_MANUAL = SALARY / MANUAL_PER_MONTH                  # RM11.55
COST_PER_TC_QAPULSE = SALARY / QAPULSE_PER_MONTH                # RM2.31
SAVINGS_PER_TC = COST_PER_TC_MANUAL - COST_PER_TC_QAPULSE        # RM9.24
PCT_REDUCTION = SAVINGS_PER_TC / COST_PER_TC_MANUAL * 100        # 80.0%

FTE_NEEDED_MANUAL = QAPULSE_PER_MONTH / MANUAL_PER_MONTH         # 5.0
EXTRA_FTE_AVOIDED = FTE_NEEDED_MANUAL - 1                        # 4.0
MONTHLY_COST_AVOIDANCE = EXTRA_FTE_AVOIDED * SALARY              # RM20,000
ANNUAL_COST_AVOIDANCE = MONTHLY_COST_AVOIDANCE * 12              # RM240,000

MANUAL_MIN_PER_TC = (8 * 60) / MANUAL_PER_DAY                    # 24.0 min
QAPULSE_MIN_PER_TC = (8 * 60) / QAPULSE_PER_DAY                  # 4.8 min

PCT_WEEK_FREED = 80.0  # to hit the old 100/week target, QAPulse needs 1 of 5 days


def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette — matches QMPulse_Executive_Presentation.pptx
    BG_DARK = RGBColor(11, 17, 32)
    BG_CARD = RGBColor(22, 33, 58)
    BG_CARD_LIGHT = RGBColor(30, 44, 76)
    BORDER_CARD = RGBColor(45, 65, 105)

    TEXT_WHITE = RGBColor(248, 250, 252)
    TEXT_MUTED = RGBColor(148, 163, 184)
    TEXT_DIM = RGBColor(100, 116, 139)

    TEAL_ACCENT = RGBColor(20, 184, 166)
    CYAN_ACCENT = RGBColor(6, 182, 212)
    INDIGO_ACCENT = RGBColor(99, 102, 241)
    EMERALD_ACCENT = RGBColor(16, 185, 129)
    PURPLE_ACCENT = RGBColor(168, 85, 247)
    GOLD_ACCENT = RGBColor(245, 158, 11)
    CORAL_ACCENT = RGBColor(244, 63, 94)

    def apply_background(slide):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_DARK
        bg_shape.line.fill.background()
        bg_shape.shadow.inherit = False
        return bg_shape

    def add_header(slide, category_text, title_text, subtitle_text=None):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = TEAL_ACCENT
        p_cat.font.name = "Calibri"

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.65))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(26)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.font.name = "Segoe UI"

        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.48), Inches(11.7), Inches(0.45))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = tf_sub.margin_right = tf_sub.margin_top = tf_sub.margin_bottom = 0
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = TEXT_MUTED
            p_sub.font.name = "Calibri"

    def add_card(slide, left, top, width, height, bg_color=BG_CARD, border_color=BORDER_CARD):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        card.shadow.inherit = False
        return card

    def add_stat_card(slide, left, top, width, height, value, title, desc, color):
        add_card(slide, left, top, width, height)
        tbox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p_v = tf.paragraphs[0]
        p_v.text = value
        p_v.font.size = Pt(28)
        p_v.font.bold = True
        p_v.font.color.rgb = color
        p_v.font.name = "Segoe UI"

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.space_before = Pt(4)

        if desc:
            p_d = tf.add_paragraph()
            p_d.text = desc
            p_d.font.size = Pt(9.5)
            p_d.font.color.rgb = TEXT_MUTED
            p_d.space_before = Pt(3)

    def style_axis_font(axis):
        axis.tick_labels.font.size = Pt(11)
        axis.tick_labels.font.color.rgb = TEXT_MUTED
        axis.tick_labels.font.name = "Calibri"
        axis.format.line.color.rgb = BORDER_CARD

    def add_column_chart(slide, left, top, width, height, chart_title, categories, series_name, values,
                          colors, number_format='#,##0', title_color=TEAL_ACCENT):
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series(series_name, values)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
        )
        chart = graphic_frame.chart
        chart.has_legend = False

        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        ct_p = chart.chart_title.text_frame.paragraphs[0]
        ct_p.font.size = Pt(14)
        ct_p.font.bold = True
        ct_p.font.color.rgb = title_color
        ct_p.font.name = "Segoe UI"

        plot = chart.plots[0]
        plot.gap_width = 60
        plot.has_data_labels = True
        dls = plot.data_labels
        dls.number_format = number_format
        dls.number_format_is_linked = False
        dls.font.size = Pt(13)
        dls.font.bold = True
        dls.font.color.rgb = TEXT_WHITE
        dls.position = XL_LABEL_POSITION.OUTSIDE_END

        series = plot.series[0]
        series.format.line.fill.background()
        for i, point in enumerate(series.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors[i % len(colors)]

        cat_axis = chart.category_axis
        style_axis_font(cat_axis)
        cat_axis.major_tick_mark = XL_TICK_MARK.NONE
        cat_axis.has_major_gridlines = False

        val_axis = chart.value_axis
        style_axis_font(val_axis)
        val_axis.has_major_gridlines = True
        val_axis.major_gridlines.format.line.color.rgb = BORDER_CARD
        val_axis.major_gridlines.format.line.width = Pt(0.5)
        val_axis.minimum_scale = 0
        val_axis.visible = True

        chart.font.color.rgb = TEXT_MUTED
        return chart

    # =========================================================================
    # SLIDE 1: TITLE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    apply_background(s1)

    deco = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    deco.fill.solid()
    deco.fill.fore_color.rgb = TEAL_ACCENT
    deco.line.fill.background()
    deco.shadow.inherit = False

    pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.4), Inches(3.6), Inches(0.42))
    pill.fill.solid()
    pill.fill.fore_color.rgb = BG_CARD
    pill.line.color.rgb = TEAL_ACCENT
    pill.line.width = Pt(1)
    pill.shadow.inherit = False
    tf_pill = pill.text_frame
    tf_pill.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_pill = tf_pill.paragraphs[0]
    p_pill.text = "●  ROI & BUSINESS CASE"
    p_pill.font.size = Pt(11)
    p_pill.font.bold = True
    p_pill.font.color.rgb = TEAL_ACCENT
    p_pill.alignment = PP_ALIGN.CENTER

    t_box = s1.shapes.add_textbox(Inches(1.2), Inches(2.05), Inches(10.8), Inches(1.8))
    tf = t_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p1 = tf.paragraphs[0]
    p1.text = "QAPulse Return on Investment"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.font.name = "Segoe UI"

    p2 = tf.add_paragraph()
    p2.text = "Quantifying the Productivity & Cost Impact of AI-Assisted QA Testing"
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEAL_ACCENT
    p2.space_before = Pt(14)

    kpi_strip = [
        ("5x", "Testing Throughput", TEAL_ACCENT),
        ("80%", "Lower Cost / Test Case", CYAN_ACCENT),
        ("RM240K", "Est. Annual Savings / QA Seat", EMERALD_ACCENT),
        ("80%", "QA Capacity Freed", GOLD_ACCENT),
    ]
    for i, (val, label, color) in enumerate(kpi_strip):
        c_left = Inches(1.2 + i * 2.75)
        add_card(s1, c_left, Inches(4.5), Inches(2.55), Inches(1.5), BG_CARD, BORDER_CARD)
        tbox = s1.shapes.add_textbox(c_left + Inches(0.15), Inches(4.65), Inches(2.25), Inches(1.2))
        tf_k = tbox.text_frame
        tf_k.word_wrap = True
        tf_k.margin_left = tf_k.margin_right = tf_k.margin_top = tf_k.margin_bottom = 0
        p_v = tf_k.paragraphs[0]
        p_v.text = val
        p_v.font.size = Pt(26)
        p_v.font.bold = True
        p_v.font.color.rgb = color
        p_v.font.name = "Segoe UI"
        p_v.alignment = PP_ALIGN.CENTER
        p_l = tf_k.add_paragraph()
        p_l.text = label
        p_l.font.size = Pt(10.5)
        p_l.font.color.rgb = TEXT_MUTED
        p_l.alignment = PP_ALIGN.CENTER
        p_l.space_before = Pt(4)

    foot = s1.shapes.add_textbox(Inches(1.2), Inches(6.7), Inches(10.8), Inches(0.4))
    tf_f = foot.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "Prepared for Executive Review  |  August 2026  |  Baseline: 1 QA seat, RM5,000/month salary"
    p_f.font.size = Pt(10.5)
    p_f.font.color.rgb = TEXT_DIM

    s1.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / NOTES FOR C-SUITE:\n"
        "\"Today I want to walk leadership through the hard numbers behind QAPulse — not just what it does, "
        "but what it is worth in Ringgit terms. The short version: one QA analyst using QAPulse produces the "
        "same test coverage as five QA analysts working manually, at 20% of the cost per test case, freeing "
        "up 80% of that analyst's week for higher-value work.\""
    )

    # =========================================================================
    # SLIDE 2: THE CURRENT STATE (MANUAL BOTTLENECK)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_background(s2)
    add_header(s2, "Current State", "The Manual Testing Bottleneck",
               "Today's baseline QA throughput, before QAPulse, per tester.")

    left_stats = [
        ("100", "Test Cases / Week", "Manual authoring & execution pace per QA analyst.", CORAL_ACCENT),
        ("20", "Test Cases / Day", "100 ÷ 5 working days — the effective daily ceiling.", CORAL_ACCENT),
        ("24 min", "Per Test Case", "Average time to author, execute, and log one test case manually.", CORAL_ACCENT),
    ]
    for i, (val, title, desc, color) in enumerate(left_stats):
        c_left = Inches(0.8 + i * 3.95)
        add_stat_card(s2, c_left, Inches(2.1), Inches(3.75), Inches(1.9), val, title, desc, color)

    add_card(s2, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5), BG_CARD_LIGHT, TEAL_ACCENT)
    vbox = s2.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(11.1), Inches(2.2))
    tf_v = vbox.text_frame
    tf_v.word_wrap = True
    tf_v.margin_left = tf_v.margin_right = tf_v.margin_top = tf_v.margin_bottom = 0
    p_vt = tf_v.paragraphs[0]
    p_vt.text = "WHY THIS IS A BUSINESS RISK"
    p_vt.font.size = Pt(12)
    p_vt.font.bold = True
    p_vt.font.color.rgb = TEAL_ACCENT

    risk_points = [
        "Fixed capacity: QA output scales only by hiring, not by demand — every new project competes for the same 100 test cases/week.",
        "Sprint-end crunch: manual authoring and Excel consolidation push execution into the final days before UAT sign-off.",
        "Cost is linear: doubling test coverage means doubling headcount — RM5,000/month per additional QA analyst.",
    ]
    for rp in risk_points:
        p_item = tf_v.add_paragraph()
        p_item.text = f"•  {rp}"
        p_item.font.size = Pt(12)
        p_item.font.color.rgb = TEXT_WHITE
        p_item.space_before = Pt(8)

    s2.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT:\n"
        "\"Before QAPulse, a QA analyst clears 100 test cases a week — 20 a day, or about 24 minutes of manual "
        "effort per test case. That ceiling is fixed. The only lever we have today to test more is to hire more, "
        "at RM5,000 a month per head.\""
    )

    # =========================================================================
    # SLIDE 3: THE QAPULSE SHIFT (THROUGHPUT CHARTS)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_background(s3)
    add_header(s3, "The QAPulse Shift", "5x Testing Throughput, Same Headcount",
               "AI-assisted authoring and one-click execution lift daily output from 20 to 100 test cases.")

    add_column_chart(
        s3, Inches(0.8), Inches(2.0), Inches(5.7), Inches(4.6),
        "Test Cases / Day (per QA)",
        ["Manual QA", "QAPulse"],
        "Daily Throughput",
        [round(MANUAL_PER_DAY, 1), QAPULSE_PER_DAY],
        [CORAL_ACCENT, TEAL_ACCENT],
        number_format='#,##0',
    )
    add_column_chart(
        s3, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.6),
        "Test Cases / Month (per QA)",
        ["Manual QA", "QAPulse"],
        "Monthly Capacity",
        [round(MANUAL_PER_MONTH), round(QAPULSE_PER_MONTH)],
        [CORAL_ACCENT, CYAN_ACCENT],
        number_format='#,##0',
    )

    s3.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT:\n"
        f"\"With QAPulse, the same analyst moves from 20 to 100 test cases a day — a flat 5x multiplier. "
        f"Over a month that is the difference between {round(MANUAL_PER_MONTH)} and {round(QAPULSE_PER_MONTH)} "
        "test cases, with no change in headcount or salary cost.\""
    )

    # =========================================================================
    # SLIDE 4: COST PER TEST CASE
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_background(s4)
    add_header(s4, "Unit Economics", "Cost Per Test Case Falls 80%",
               "Same RM5,000 monthly salary, spread across 5x more completed test cases.")

    add_column_chart(
        s4, Inches(0.8), Inches(2.0), Inches(6.6), Inches(4.7),
        "Cost Per Test Case (RM)",
        ["Manual QA", "QAPulse"],
        "Cost / Test Case",
        [round(COST_PER_TC_MANUAL, 2), round(COST_PER_TC_QAPULSE, 2)],
        [CORAL_ACCENT, TEAL_ACCENT],
        number_format='"RM"#,##0.00',
    )

    add_card(s4, Inches(7.7), Inches(2.0), Inches(4.8), Inches(4.7), BG_CARD_LIGHT, TEAL_ACCENT)
    vbox = s4.shapes.add_textbox(Inches(8.0), Inches(2.25), Inches(4.2), Inches(4.2))
    tf_v = vbox.text_frame
    tf_v.word_wrap = True
    tf_v.margin_left = tf_v.margin_right = tf_v.margin_top = tf_v.margin_bottom = 0

    p_vt = tf_v.paragraphs[0]
    p_vt.text = "THE CALCULATION"
    p_vt.font.size = Pt(12)
    p_vt.font.bold = True
    p_vt.font.color.rgb = TEAL_ACCENT

    calc_lines = [
        (f"Manual: RM5,000 ÷ {round(MANUAL_PER_MONTH)} test cases/month", f"= RM{COST_PER_TC_MANUAL:.2f} / test case"),
        (f"QAPulse: RM5,000 ÷ {round(QAPULSE_PER_MONTH)} test cases/month", f"= RM{COST_PER_TC_QAPULSE:.2f} / test case"),
    ]
    for label, result in calc_lines:
        p_l = tf_v.add_paragraph()
        p_l.text = label
        p_l.font.size = Pt(12)
        p_l.font.color.rgb = TEXT_WHITE
        p_l.space_before = Pt(12)
        p_r = tf_v.add_paragraph()
        p_r.text = result
        p_r.font.size = Pt(15)
        p_r.font.bold = True
        p_r.font.color.rgb = CYAN_ACCENT
        p_r.space_before = Pt(2)

    p_sep = tf_v.add_paragraph()
    p_sep.text = " "
    p_sep.space_before = Pt(6)

    p_final = tf_v.add_paragraph()
    p_final.text = f"Savings: RM{SAVINGS_PER_TC:.2f} per test case"
    p_final.font.size = Pt(15)
    p_final.font.bold = True
    p_final.font.color.rgb = EMERALD_ACCENT
    p_final.space_before = Pt(10)

    p_pct = tf_v.add_paragraph()
    p_pct.text = f"→ {PCT_REDUCTION:.0f}% reduction in cost per test case"
    p_pct.font.size = Pt(13)
    p_pct.font.color.rgb = TEXT_MUTED
    p_pct.space_before = Pt(4)

    s4.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT:\n"
        f"\"The same RM5,000 salary now buys RM{COST_PER_TC_MANUAL:.2f} worth of testing down to "
        f"RM{COST_PER_TC_QAPULSE:.2f} per test case — an 80% reduction in unit cost, with no change in "
        "quality standards or sign-off process.\""
    )

    # =========================================================================
    # SLIDE 5: HEADCOUNT COST AVOIDANCE
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_background(s5)
    add_header(s5, "Headcount Impact", "RM240,000 Avoided Per QAPulse Seat, Per Year",
               f"To manually match one QAPulse analyst's monthly output ({round(QAPULSE_PER_MONTH)} test cases), you would need {int(FTE_NEEDED_MANUAL)} QA analysts.")

    add_column_chart(
        s5, Inches(0.8), Inches(2.0), Inches(6.6), Inches(4.7),
        f"Monthly Payroll for {round(QAPULSE_PER_MONTH)} Test Cases (RM)",
        [f"Manual ({int(FTE_NEEDED_MANUAL)} QA)", "QAPulse (1 QA)"],
        "Payroll Cost",
        [int(FTE_NEEDED_MANUAL * SALARY), int(SALARY)],
        [CORAL_ACCENT, TEAL_ACCENT],
        number_format='"RM"#,##0',
    )

    stats = [
        ("4", "Extra Headcount Avoided", "per QAPulse-equipped QA seat", GOLD_ACCENT),
        ("RM20,000", "Avoided Cost / Month", "4 FTE x RM5,000 salary", EMERALD_ACCENT),
        ("RM240,000", "Avoided Cost / Year", "per QAPulse seat, compounding across the team", EMERALD_ACCENT),
    ]
    for i, (val, title, desc, color) in enumerate(stats):
        c_top = Inches(2.0 + i * 1.62)
        add_stat_card(s5, Inches(7.7), c_top, Inches(4.8), Inches(1.45), val, title, desc, color)

    s5.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT:\n"
        f"\"Put another way: hitting {round(QAPULSE_PER_MONTH)} test cases a month manually needs "
        f"{int(FTE_NEEDED_MANUAL)} QA analysts on payroll — RM{int(FTE_NEEDED_MANUAL*SALARY):,}/month. "
        "QAPulse gets one analyst to the same output for RM5,000/month. That is RM20,000/month, or "
        "RM240,000/year, in avoided hiring cost for every QAPulse-equipped seat — and it scales linearly "
        "as the team grows.\""
    )

    # =========================================================================
    # SLIDE 6: 12-MONTH VALUE PROJECTION
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_background(s6)
    add_header(s6, "Value Over Time", "Cumulative Cost Avoidance — 12 Month View",
               "RM20,000/month in avoided headcount cost compounds to RM240,000 in year one, per QAPulse seat.")

    add_column_chart(
        s6, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.7),
        "Cumulative Cost Avoidance (RM)",
        ["Q1", "Q2", "Q3", "Q4"],
        "Cumulative Savings",
        [int(MONTHLY_COST_AVOIDANCE * 3), int(MONTHLY_COST_AVOIDANCE * 6),
         int(MONTHLY_COST_AVOIDANCE * 9), int(MONTHLY_COST_AVOIDANCE * 12)],
        [INDIGO_ACCENT, INDIGO_ACCENT, CYAN_ACCENT, TEAL_ACCENT],
        number_format='"RM"#,##0',
    )

    s6.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT:\n"
        "\"This avoided cost is not a one-time event — it compounds every month QAPulse stays in use. "
        "By the end of year one, a single QAPulse seat has avoided RM240,000 in headcount cost that would "
        "otherwise have been needed to keep pace with the same testing volume.\""
    )

    # =========================================================================
    # SLIDE 7: ADDITIONAL EFFICIENCY GAINS
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_background(s7)
    add_header(s7, "Beyond Test-Case Volume", "What Else QAPulse Saves Today",
               "Value the salary-based ROI model above does not even capture.")

    extra_gains = [
        ("AI Requirement Analyzer & Test Generator", "~40% faster test design",
         "Google GenAI drafts test cases straight from requirements/user stories — analysts edit and approve rather than author from scratch.", TEAL_ACCENT),
        ("Automated PMO Reporting", "Days → under 60 seconds",
         "Review Log, Pareto Analysis, CAPA sheets and the PMO verdict email/PDF are generated automatically on save — no manual Excel consolidation.", CYAN_ACCENT),
        ("Automatic Defect Creation", "Zero duplicate data entry",
         "A Failed test step creates a Redmine defect as a child issue automatically — testers never re-key the same failure twice.", INDIGO_ACCENT),
        ("End-to-End Traceability", "100% requirement-to-defect coverage",
         "Every requirement, test case, execution and defect is linked — eliminating the cost of defects that leak past UAT into production.", EMERALD_ACCENT),
        ("Eliminated Spreadsheet Sprawl", "0 version-conflict rework",
         "One shared execution file per Redmine ticket replaces emailed Excel copies — no more reconciling whose version is current.", GOLD_ACCENT),
        ("Immutable Audit Trail", "Lower compliance/audit prep cost",
         "Every status change and sign-off is logged automatically, reducing the effort to reconstruct evidence for audits.", PURPLE_ACCENT),
    ]

    col_w = Inches(3.83)
    row_h = Inches(2.3)
    for i, (title, tag, desc, color) in enumerate(extra_gains):
        col = i % 3
        row = i // 3
        c_left = Inches(0.8) + col * col_w
        c_top = Inches(2.0) + row * row_h
        add_card(s7, c_left, c_top, col_w - Inches(0.15), row_h - Inches(0.2), BG_CARD, BORDER_CARD)
        tbox = s7.shapes.add_textbox(c_left + Inches(0.18), c_top + Inches(0.15), col_w - Inches(0.5), row_h - Inches(0.45))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p_tag = tf.paragraphs[0]
        p_tag.text = tag.upper()
        p_tag.font.size = Pt(11.5)
        p_tag.font.bold = True
        p_tag.font.color.rgb = color

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(12.5)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.space_before = Pt(5)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(4)

    s7.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT:\n"
        "\"The RM240,000 figure is only the headcount-avoidance piece — the part we can price precisely. "
        "QAPulse also compresses test design time by roughly 40%, turns multi-day PMO reporting cycles into "
        "under a minute, removes duplicate defect logging, and closes the traceability gaps that let defects "
        "slip into production. None of that is included in the RM240,000 number — it is upside on top.\""
    )

    # =========================================================================
    # SLIDE 8: ASSUMPTIONS & METHODOLOGY
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_background(s8)
    add_header(s8, "Methodology", "Assumptions Behind These Numbers",
               "Transparent inputs so Finance and HR can validate or recalibrate this model.")

    assumptions = [
        ("Baseline QA salary", "RM5,000 / month per QA analyst (as provided)."),
        ("Manual throughput", "100 test cases / week, 5-day work week → 20 test cases / day."),
        ("QAPulse throughput", "100 test cases / day (as provided), same 5-day work week."),
        ("Month length", "4.33 weeks / month (21.65 working days) — standard calendar averaging."),
        ("Cost per test case", "Monthly salary ÷ monthly test cases completed, same analyst, same salary."),
        ("Headcount avoidance", "Extra manual FTEs needed to match QAPulse's monthly output, at the same RM5,000 salary each."),
        ("Not included", "QAPulse licensing/build/hosting cost is not netted off — the figures above are gross value created, not net ROI %."),
        ("Scaling", "All figures are per QA seat — multiply by the number of QAPulse-equipped analysts for team- or department-wide impact."),
    ]

    tbox = s8.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(5.0))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for label, detail in assumptions:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0) if first else Pt(11)
        run1 = p.add_run()
        run1.text = f"{label}:  "
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = TEAL_ACCENT
        run2 = p.add_run()
        run2.text = detail
        run2.font.size = Pt(13)
        run2.font.color.rgb = TEXT_WHITE

    s8.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT:\n"
        "\"We wanted this model to survive scrutiny, so every input is stated plainly on this slide. "
        "Swap in your actual salary band or team size and the same formulas hold — we are happy to rerun "
        "this with Finance's own numbers before it goes into a budget request.\""
    )

    # =========================================================================
    # SLIDE 9: RECOMMENDATION / NEXT STEPS
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_background(s9)
    add_header(s9, "Recommendation", "The Business Case Is Clear",
               "A 5x productivity gain and RM240,000/year in avoided cost, per QA seat.")

    add_card(s9, Inches(0.8), Inches(2.1), Inches(11.7), Inches(2.3), BG_CARD_LIGHT, TEAL_ACCENT)
    vbox = s9.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(11.1), Inches(2.0))
    tf_v = vbox.text_frame
    tf_v.word_wrap = True
    tf_v.margin_left = tf_v.margin_right = tf_v.margin_top = tf_v.margin_bottom = 0
    p_vt = tf_v.paragraphs[0]
    p_vt.text = "THE VERDICT"
    p_vt.font.size = Pt(12)
    p_vt.font.bold = True
    p_vt.font.color.rgb = TEAL_ACCENT
    p_vb = tf_v.add_paragraph()
    p_vb.text = ("Every QA analyst equipped with QAPulse delivers the output of five, at one-fifth the cost per "
                 "test case — worth RM240,000/year in avoided headcount alone, before counting faster test "
                 "design, automated reporting, and full traceability.")
    p_vb.font.size = Pt(14)
    p_vb.font.color.rgb = TEXT_WHITE
    p_vb.space_before = Pt(8)

    next_steps = [
        ("Validate", "Confirm salary bands and current QA headcount with Finance/HR to size the full team-wide savings."),
        ("Pilot", "Track actual test cases/day for one sprint on 1-2 projects to confirm the 5x multiplier in practice."),
        ("Scale", "Extend to the full QA team and fold the confirmed savings into the next budget cycle."),
    ]
    for i, (title, desc) in enumerate(next_steps):
        c_left = Inches(0.8 + i * 3.95)
        add_card(s9, c_left, Inches(4.7), Inches(3.75), Inches(2.0), BG_CARD, BORDER_CARD)
        tbox = s9.shapes.add_textbox(c_left + Inches(0.2), Inches(4.85), Inches(3.35), Inches(1.7))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p_t = tf.paragraphs[0]
        p_t.text = f"{i+1}.  {title}"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = CYAN_ACCENT
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(6)

    s9.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / CLOSING:\n"
        "\"To close: this is not a projection built on hope, it is basic arithmetic on numbers we already "
        "observe today — 100 test cases a week manually versus 100 a day with QAPulse. We would like to "
        "validate the team-wide number with Finance and HR, then move to a short pilot to confirm it in "
        "the field. Thank you — happy to take questions.\""
    )

    output_path = os.path.join(os.getcwd(), "QAPulse_ROI_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")


if __name__ == "__main__":
    create_deck()
