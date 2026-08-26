import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 Widescreen standard
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette: Modern Executive Dark Slate & Electric Teal / Cyan
    BG_DARK = RGBColor(11, 17, 32)        # Deep Obsidian Slate #0B1120
    BG_CARD = RGBColor(22, 33, 58)        # Elevated Slate Card #16213A
    BG_CARD_LIGHT = RGBColor(30, 44, 76)  # Accent Card #1E2C4C
    BORDER_CARD = RGBColor(45, 65, 105)   # Border #2D4169
    
    TEXT_WHITE = RGBColor(248, 250, 252)  # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8
    TEXT_DIM = RGBColor(100, 116, 139)    # #64748B
    
    TEAL_ACCENT = RGBColor(20, 184, 166)  # Electric Teal #14B8A6
    CYAN_ACCENT = RGBColor(6, 182, 212)   # Bright Cyan #06B6D4
    INDIGO_ACCENT = RGBColor(99, 102, 241)# Indigo Accent #6366F1
    EMERALD_ACCENT = RGBColor(16, 185, 129)# Emerald #10B981
    PURPLE_ACCENT = RGBColor(168, 85, 247)# Purple #A855F7
    GOLD_ACCENT = RGBColor(245, 158, 11)  # Amber Gold #F59E0B
    CORAL_ACCENT = RGBColor(244, 63, 94)  # Coral Red #F43F5E

    def apply_background(slide):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_DARK
        bg_shape.line.fill.background()
        return bg_shape

    def add_header(slide, category_text, title_text, subtitle_text=None):
        # Category / Pill
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = TEAL_ACCENT
        p_cat.font.name = "Calibri"

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.65))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(26)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.font.name = "Segoe UI"

        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.48), Inches(11.7), Inches(0.45))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = tf_sub.margin_right = tf_sub.margin_top = tf_sub.margin_bottom = 0
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = TEXT_MUTED
            p_sub.font.name = "Calibri"

    def add_card(slide, left, top, width, height, bg_color=BG_CARD, border_color=BORDER_CARD):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    apply_background(s1)

    # Accent decorative top bar
    deco = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    deco.fill.solid()
    deco.fill.fore_color.rgb = TEAL_ACCENT
    deco.line.fill.background()

    # Category Pill
    pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.5), Inches(3.2), Inches(0.42))
    pill.fill.solid()
    pill.fill.fore_color.rgb = BG_CARD
    pill.line.color.rgb = TEAL_ACCENT
    pill.line.width = Pt(1)
    tf_pill = pill.text_frame
    tf_pill.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_pill = tf_pill.paragraphs[0]
    p_pill.text = "●  EXECUTIVE STRATEGY PITCH"
    p_pill.font.size = Pt(11)
    p_pill.font.bold = True
    p_pill.font.color.rgb = TEAL_ACCENT
    p_pill.alignment = PP_ALIGN.CENTER

    # Main Hero Title
    t_box = s1.shapes.add_textbox(Inches(1.2), Inches(2.15), Inches(10.8), Inches(1.8))
    tf = t_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p1 = tf.paragraphs[0]
    p1.text = "QMPulse (Quality Management Pulse)"
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.font.name = "Segoe UI"

    p2 = tf.add_paragraph()
    p2.text = "From First Requirement to Final Verdict — One Living Heartbeat."
    p2.font.size = Pt(22)
    p2.font.color.rgb = CYAN_ACCENT
    p2.font.name = "Segoe UI"
    p2.space_before = Pt(8)

    # Description
    desc_box = s1.shapes.add_textbox(Inches(1.2), Inches(4.1), Inches(10.8), Inches(1.0))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = tf_desc.margin_top = tf_desc.margin_right = tf_desc.margin_bottom = 0
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "A Next-Generation Delivery Quality & Governance Platform uniting PMO, Functional Analysts, Developers, and QA into a single, automated, Redmine-synced and AI-accelerated ecosystem."
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.font.name = "Calibri"

    # Stat Highlights Footer Grid (4 Key Numbers)
    stat_items = [
        ("100%", "End-to-End Traceability"),
        ("8 Stages", "Unified Milestone Lifecycle"),
        ("15 Roles", "Strict RBAC Governance"),
        ("1 Source", "Single Truth (Redmine + AI)")
    ]
    for i, (val, lbl) in enumerate(stat_items):
        c_left = Inches(1.2 + i * 2.8)
        add_card(s1, c_left, Inches(5.4), Inches(2.6), Inches(1.2), BG_CARD, BORDER_CARD)
        st_box = s1.shapes.add_textbox(c_left + Inches(0.15), Inches(5.5), Inches(2.3), Inches(1.0))
        tf_st = st_box.text_frame
        tf_st.word_wrap = True
        tf_st.margin_left = tf_st.margin_right = tf_st.margin_top = tf_st.margin_bottom = 0
        p_v = tf_st.paragraphs[0]
        p_v.text = val
        p_v.font.size = Pt(20)
        p_v.font.bold = True
        p_v.font.color.rgb = TEAL_ACCENT
        p_v.font.name = "Segoe UI"
        
        p_l = tf_st.add_paragraph()
        p_l.text = lbl
        p_l.font.size = Pt(11)
        p_l.font.color.rgb = TEXT_MUTED
        p_l.font.name = "Calibri"
        p_l.space_before = Pt(3)

    s1.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / NOTES FOR C-SUITE:\n"
        "\"Good morning, Executive Leadership. Today, I am proud to present QMPulse—our unified Quality Management platform.\n"
        "In software and product delivery, quality cannot simply be an afterthought or an isolated QA phase before release. "
        "QMPulse connects the entire delivery chain: from the very first business requirement authored by Functional Analysts, "
        "through sprint build, QA test execution, and real-time PMO governance, straight to final go-live sign-off. "
        "It gives leadership 100% end-to-end visibility, guarantees compliance, eliminates spreadsheet chaos, and accelerates testing with built-in AI.\""
    )

    # =========================================================================
    # SLIDE 2: THE STRATEGIC PROBLEM & THE QMPULSE SOLUTION
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_background(s2)
    add_header(s2, "Strategic Context", "The Delivery Dilemma vs. The Unified Pulse", "Why traditional enterprise delivery breaks down—and how QMPulse fixes it.")

    # Left Column: The Problem (Current Friction)
    add_card(s2, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8), BG_CARD, RGBColor(180, 50, 50))
    p_badge = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(2.2), Inches(2.2), Inches(0.35))
    p_badge.fill.solid()
    p_badge.fill.fore_color.rgb = RGBColor(60, 20, 25)
    p_badge.line.color.rgb = CORAL_ACCENT
    tf_pb = p_badge.text_frame
    tf_pb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_pbt = tf_pb.paragraphs[0]
    p_pbt.text = "CURRENT CHALLENGES"
    p_pbt.font.size = Pt(10)
    p_pbt.font.bold = True
    p_pbt.font.color.rgb = CORAL_ACCENT
    p_pbt.alignment = PP_ALIGN.CENTER

    prob_box = s2.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(5.0), Inches(3.9))
    tf_prob = prob_box.text_frame
    tf_prob.word_wrap = True
    tf_prob.margin_left = tf_prob.margin_right = tf_prob.margin_top = tf_prob.margin_bottom = 0

    problems = [
        ("Fragmented Silos", "PMO, Functional Analysts, Devs, and QA work in disconnected tools, creating communication latency and missed handoffs."),
        ("Spreadsheet & Email Chaos", "Test cases, execution trackers, and defect logs get passed around on Excel attachments with zero version control."),
        ("Untracked Defect Drift", "Defects discovered during testing are disconnected from Redmine tickets, causing rework loops and untracked scope."),
        ("Blind Executive Oversight", "PMO & Leadership lack real-time SPI, first-pass yield, and release readiness data until it is too late to react.")
    ]
    for i, (title, body) in enumerate(problems):
        p_t = tf_prob.add_paragraph() if i > 0 else tf_prob.paragraphs[0]
        p_t.text = f"✕  {title}"
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = RGBColor(251, 113, 133)
        if i > 0: p_t.space_before = Pt(10)
        
        p_b = tf_prob.add_paragraph()
        p_b.text = body
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(2)

    # Right Column: The Solution (QMPulse Advantage)
    add_card(s2, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.8), BG_CARD, RGBColor(20, 184, 166))
    s_badge = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(2.2), Inches(2.2), Inches(0.35))
    s_badge.fill.solid()
    s_badge.fill.fore_color.rgb = RGBColor(10, 45, 45)
    s_badge.line.color.rgb = TEAL_ACCENT
    tf_sb = s_badge.text_frame
    tf_sb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_sbt = tf_sb.paragraphs[0]
    p_sbt.text = "THE QMPULSE SOLUTION"
    p_sbt.font.size = Pt(10)
    p_sbt.font.bold = True
    p_sbt.font.color.rgb = TEAL_ACCENT
    p_sbt.alignment = PP_ALIGN.CENTER

    sol_box = s2.shapes.add_textbox(Inches(7.1), Inches(2.7), Inches(5.1), Inches(3.9))
    tf_sol = sol_box.text_frame
    tf_sol.word_wrap = True
    tf_sol.margin_left = tf_sol.margin_right = tf_sol.margin_top = tf_sol.margin_bottom = 0

    solutions = [
        ("Unified Living Heartbeat", "Single shared system of record orchestrating PMO, FA, Dev, QA, and Leadership in real-time continuous sync."),
        ("Native Two-Way Redmine Sync", "Integrates directly with enterprise Redmine; defects, requirements, and progress synchronize bi-directionally without disruption."),
        ("AI-Accelerated Quality Hub", "Google GenAI analyzes requirements for ambiguities and auto-generates comprehensive test case matrices in seconds."),
        ("Executive Single Pane of Glass", "Instant live metrics on Schedule Performance Index (SPI), stability, first-pass yield, and automated PMO verdict reports.")
    ]
    for i, (title, body) in enumerate(solutions):
        p_t = tf_sol.add_paragraph() if i > 0 else tf_sol.paragraphs[0]
        p_t.text = f"✓  {title}"
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = TEAL_ACCENT
        if i > 0: p_t.space_before = Pt(10)
        
        p_b = tf_sol.add_paragraph()
        p_b.text = body
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(2)

    s2.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / NOTES FOR C-SUITE:\n"
        "\"To understand why QMPulse is vital, let us look at the industry's common bottleneck: tool sprawl.\n"
        "Currently, PMOs track in one tool, developers write in Redmine/Jira, FAs write specs in Word or Confluence, and QA runs test cases on spreadsheets. "
        "This leads to version skew, untracked defects, and eleventh-hour release panic. \n"
        "QMPulse solves this by establishing one single operational pulse. We are not discarding Redmine—we are supercharging it with automated workflows, "
        "AI assistance, and instant executive analytics. Everything is transparent, tracked, and verifiable.\""
    )

    # =========================================================================
    # SLIDE 3: FIVE TEAMS, ONE LIVING HEARTBEAT
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_background(s3)
    add_header(s3, "Stakeholder Alignment", "Five Delivery Teams. One Synchronized Cadence.", "Tailored operational views for every discipline on the same real-time data.")

    personas = [
        ("Project & PMO", "Own the Milestone", PURPLE_ACCENT, [
            "Milestone & sprint planning",
            "Environment & capacity allocation",
            "Live schedule health & SPI tracking",
            "Continuous risk & mitigation log"
        ]),
        ("Functional Analysts", "Sharper Requirements", INDIGO_ACCENT, [
            "Structured requirement authoring",
            "AI Requirement Analyzer (gap detection)",
            "Strict peer-review & segregation",
            "Zero self-approvals enforced"
        ]),
        ("Developers", "Build in Context", CYAN_ACCENT, [
            "Pick up approved specs directly",
            "Build with full requirement context",
            "1-click status flip to 'For QA Test'",
            "Direct Redmine issue sync back"
        ]),
        ("QA Engineers", "Prove Every Path", TEAL_ACCENT, [
            "Author test cases during build phase",
            "AI-assisted test suite generation",
            "Environment-specific test runs",
            "Auto-raise Redmine defects on failure"
        ]),
        ("C-Suite & Leadership", "The Pulse at a Glance", EMERALD_ACCENT, [
            "Single-screen executive dashboard",
            "First-pass rate & defect stability",
            "Automated PMO verdict emails/PDFs",
            "Audit-ready compliance & Pareto logs"
        ]),
    ]

    for i, (role, tagline, color, points) in enumerate(personas):
        c_left = Inches(0.8 + i * 2.4)
        add_card(s3, c_left, Inches(2.1), Inches(2.25), Inches(4.7), BG_CARD, BORDER_CARD)

        # Top Accent Header inside card
        top_bar = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(2.1), Inches(2.25), Inches(0.7))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = BG_CARD_LIGHT
        top_bar.line.color.rgb = color
        top_bar.line.width = Pt(1)
        tf_tb = top_bar.text_frame
        tf_tb.word_wrap = True
        tf_tb.margin_left = tf_tb.margin_right = Inches(0.08)
        p_role = tf_tb.paragraphs[0]
        p_role.text = role.upper()
        p_role.font.size = Pt(10)
        p_role.font.bold = True
        p_role.font.color.rgb = color
        p_role.font.name = "Calibri"

        # Content Box
        tbox = s3.shapes.add_textbox(c_left + Inches(0.1), Inches(2.9), Inches(2.05), Inches(3.8))
        tf_c = tbox.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = 0

        p_tag = tf_c.paragraphs[0]
        p_tag.text = tagline
        p_tag.font.size = Pt(13)
        p_tag.font.bold = True
        p_tag.font.color.rgb = TEXT_WHITE
        p_tag.font.name = "Segoe UI"

        for pt in points:
            p_pt = tf_c.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.size = Pt(10)
            p_pt.font.color.rgb = TEXT_MUTED
            p_pt.font.name = "Calibri"
            p_pt.space_before = Pt(6)

    s3.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / NOTES FOR C-SUITE:\n"
        "\"Here is how QMPulse serves every single leader in this room:\n"
        "- For our PMO: You open the milestone, assign capacity, track Schedule Performance Index (SPI), and manage risks.\n"
        "- For Functional Analysts: Requirements are captured with AI gap analysis, and peer-approval enforces segregation of duties.\n"
        "- For Dev Leads & Developers: Tasks flow seamlessly from approved requirements to code with zero ambiguity.\n"
        "- For QA Teams: Test cases are pre-authored while dev is coding. The moment a build lands, test execution begins.\n"
        "- For the C-Suite: You get the macro-level pulse: release stability, first-pass yield, and automated audit trails.\""
    )

    # =========================================================================
    # SLIDE 4: THE 8-BEAT MILESTONE LIFECYCLE
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_background(s4)
    add_header(s4, "End-to-End Workflow", "The 8 Beats of a Milestone Lifecycle", "From intake to go-live sign-off with enforced segregation of duties and automated gates.")

    stages_row1 = [
        ("01", "Plan Milestone", "PMO", PURPLE_ACCENT, "Define milestone scope, target dates, environments & capacity allocation."),
        ("02", "Author Specs", "FA + AI", PURPLE_ACCENT, "Capture requirements; AI Analyzer flags ambiguities and untestable logic."),
        ("03", "Peer-Approve", "FA Reviewer", INDIGO_ACCENT, "Peer review enforced; strict segregation prevents self-sign-offs."),
        ("04", "Assign & Build", "Dev Lead / Dev", INDIGO_ACCENT, "Dev builds against approved specs; flips ticket to 'For QA Test'.")
    ]

    stages_row2 = [
        ("05", "Author Tests", "QA + AI", CYAN_ACCENT, "Test cases authored in parallel with build; AI generates full edge cases."),
        ("06", "Execute Testing", "QA Engineers", TEAL_ACCENT, "Environment-level execution; failing steps auto-raise Redmine child defects."),
        ("07", "Track & Close", "QA Lead", TEAL_ACCENT, "Real-time pass/fail/blocked analytics and blocker resolution."),
        ("08", "UAT & Go-Live", "UAT & PMO", EMERALD_ACCENT, "Business UAT sign-off, milestone close, and automated lessons learned capture.")
    ]

    for row_idx, row_stages in enumerate([stages_row1, stages_row2]):
        top_y = Inches(2.1 + row_idx * 2.3)
        for col_idx, (num, title, owner, color, desc) in enumerate(row_stages):
            c_left = Inches(0.8 + col_idx * 2.95)
            add_card(s4, c_left, top_y, Inches(2.8), Inches(2.1), BG_CARD, BORDER_CARD)

            # Stage Number Badge
            num_box = s4.shapes.add_shape(MSO_SHAPE.OVAL, c_left + Inches(0.15), top_y + Inches(0.15), Inches(0.42), Inches(0.42))
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = color
            num_box.line.fill.background()
            tf_n = num_box.text_frame
            tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_n = tf_n.paragraphs[0]
            p_n.text = num
            p_n.font.size = Pt(11)
            p_n.font.bold = True
            p_n.font.color.rgb = BG_DARK
            p_n.alignment = PP_ALIGN.CENTER

            # Text content
            sbox = s4.shapes.add_textbox(c_left + Inches(0.65), top_y + Inches(0.12), Inches(2.05), Inches(1.85))
            tf_s = sbox.text_frame
            tf_s.word_wrap = True
            tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = tf_s.margin_bottom = 0

            p_t = tf_s.paragraphs[0]
            p_t.text = title
            p_t.font.size = Pt(13)
            p_t.font.bold = True
            p_t.font.color.rgb = TEXT_WHITE
            p_t.font.name = "Segoe UI"

            p_o = tf_s.add_paragraph()
            p_o.text = f"Owner: {owner}"
            p_o.font.size = Pt(10)
            p_o.font.bold = True
            p_o.font.color.rgb = color
            p_o.font.name = "Calibri"
            p_o.space_before = Pt(2)

            p_d = tf_s.add_paragraph()
            p_d.text = desc
            p_d.font.size = Pt(10)
            p_d.font.color.rgb = TEXT_MUTED
            p_d.font.name = "Calibri"
            p_d.space_before = Pt(4)

    # Footer banner on slide 4: Cross-cutting pillars
    banner = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.75), Inches(11.7), Inches(0.42))
    banner.fill.solid()
    banner.fill.fore_color.rgb = BG_CARD_LIGHT
    banner.line.color.rgb = TEAL_ACCENT
    banner.line.width = Pt(1)
    tf_b = banner.text_frame
    tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_b = tf_b.paragraphs[0]
    p_b.text = "⚡ RUNNING THROUGH EVERY BEAT:  Instant Smart Notifications  |  Dynamic Risk Register  |  Continuous Lessons Learned"
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.color.rgb = TEAL_ACCENT
    p_b.alignment = PP_ALIGN.CENTER

    s4.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / NOTES FOR C-SUITE:\n"
        "\"This 8-stage lifecycle is the core engine of QMPulse. \n"
        "Notice the built-in governance gates: \n"
        "1. Segregation of Duties: A Functional Analyst can never sign off their own requirement—it must be peer-reviewed by a colleague in the same project.\n"
        "2. Parallel Acceleration: QA writes test cases during the build phase (Stage 5) so no time is wasted when code is deployed.\n"
        "3. Automated Defect Linking: When QA fails a step in Stage 6, a defect is automatically created in Redmine as a child ticket of the original requirement.\n"
        "4. Risk Management is active through all 8 stages, not just in stage 1.\""
    )

    # =========================================================================
    # SLIDE 5: THE ENTERPRISE BACKBONE: REDMINE + AI HUB
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_background(s5)
    add_header(s5, "Technology Backbone", "Deep Redmine Integration & The Enterprise AI Hub", "Zero disruption to existing systems of record, multiplied by intelligent GenAI assistance.")

    # Left Column: Redmine Two-Way Sync
    add_card(s5, Inches(0.8), Inches(2.0), Inches(5.65), Inches(4.8), BG_CARD, INDIGO_ACCENT)
    
    r_badge = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(2.2), Inches(2.4), Inches(0.35))
    r_badge.fill.solid()
    r_badge.fill.fore_color.rgb = RGBColor(30, 27, 75)
    r_badge.line.color.rgb = INDIGO_ACCENT
    tf_rb = r_badge.text_frame
    tf_rb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_rbt = tf_rb.paragraphs[0]
    p_rbt.text = "DEEP REDMINE INTEGRATION"
    p_rbt.font.size = Pt(10)
    p_rbt.font.bold = True
    p_rbt.font.color.rgb = RGBColor(165, 180, 252)
    p_rbt.alignment = PP_ALIGN.CENTER

    r_box = s5.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(5.05), Inches(3.9))
    tf_r = r_box.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0

    redmine_pts = [
        ("Non-Disruptive Adoption", "QMPulse leverages your existing Redmine installation—no costly data migrations or team re-training required."),
        ("Two-Way Real-Time Synchronization", "Requirements and tickets import with full hierarchy; defects and status transitions write straight back."),
        ("Safe Insert-Only Sync Architecture", "Existing Redmine data is never overwritten or clobbered; project custom fields map seamlessly."),
        ("Auto-Defect Creation on Test Failure", "One click converts a failed test step into a fully populated Redmine defect with logs and attachments.")
    ]
    for i, (title, body) in enumerate(redmine_pts):
        p_t = tf_r.add_paragraph() if i > 0 else tf_r.paragraphs[0]
        p_t.text = f"◆  {title}"
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = RGBColor(165, 180, 252)
        if i > 0: p_t.space_before = Pt(8)
        
        p_b = tf_r.add_paragraph()
        p_b.text = body
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(2)

    # Right Column: The AI Hub
    add_card(s5, Inches(6.85), Inches(2.0), Inches(5.65), Inches(4.8), BG_CARD, TEAL_ACCENT)
    
    ai_badge = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.15), Inches(2.2), Inches(2.4), Inches(0.35))
    ai_badge.fill.solid()
    ai_badge.fill.fore_color.rgb = RGBColor(6, 40, 45)
    ai_badge.line.color.rgb = TEAL_ACCENT
    tf_ab = ai_badge.text_frame
    tf_ab.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_abt = tf_ab.paragraphs[0]
    p_abt.text = "THE ENTERPRISE AI HUB"
    p_abt.font.size = Pt(10)
    p_abt.font.bold = True
    p_abt.font.color.rgb = TEAL_ACCENT
    p_abt.alignment = PP_ALIGN.CENTER

    ai_box = s5.shapes.add_textbox(Inches(7.15), Inches(2.7), Inches(5.05), Inches(3.9))
    tf_ai = ai_box.text_frame
    tf_ai.word_wrap = True
    tf_ai.margin_left = tf_ai.margin_right = tf_ai.margin_top = tf_ai.margin_bottom = 0

    ai_pts = [
        ("AI Requirement Analyzer", "Evaluates functional specs before development starts; flags ambiguities, missing edge cases, and untestable logic."),
        ("Automated Test Case Generation", "Generates comprehensive, multi-step test cases from user stories with expected results and boundary values."),
        ("QA Copilot Workspace Assistant", "Natural-language assistant answering queries across test suites, execution history, and delivery metrics."),
        ("Requirement Chat & Refinement", "Conversational requirement refinement allowing FAs and QA to challenge and perfect specifications.")
    ]
    for i, (title, body) in enumerate(ai_pts):
        p_t = tf_ai.add_paragraph() if i > 0 else tf_ai.paragraphs[0]
        p_t.text = f"✦  {title}"
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = TEAL_ACCENT
        if i > 0: p_t.space_before = Pt(8)
        
        p_b = tf_ai.add_paragraph()
        p_b.text = body
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(2)

    s5.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / NOTES FOR C-SUITE:\n"
        "\"Two strategic technological pillars set QMPulse apart:\n"
        "1. We integrate seamlessly with Redmine. Many platforms force you to migrate data or rip out existing tracking. "
        "QMPulse acts as a modern quality layer over Redmine, synchronizing in real time without overwriting history.\n"
        "2. We embed Google GenAI directly into the workflow. Instead of using AI as a gimmick, our AI Requirement Analyzer acts as an early quality gate "
        "to catch bad requirements before they reach development, while AI Test Generation cuts test design effort by over 40%.\""
    )

    # =========================================================================
    # SLIDE 6: COMPREHENSIVE CAPABILITY MATRIX (4 DISCIPLINES)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_background(s6)
    add_header(s6, "Platform Breadth", "Enterprise Capabilities Across 4 Core Disciplines", "12 integrated modules delivering an end-to-end quality and governance command center.")

    pillars = [
        ("PLAN & GOVERN", PURPLE_ACCENT, [
            ("Milestones & Sprints", "Scope, environment targets, and real-time schedule health."),
            ("Resource Allocation", "Capacity visibility to spot team bottlenecks before they cause delays."),
            ("Risk Register", "Dynamic risk scoring, mitigation plans, and ownership assignment.")
        ]),
        ("BUILD & TEST", CYAN_ACCENT, [
            ("Test Case Library", "Versioned test cases with steps, expected outcomes & module mapping."),
            ("Execution Dashboard", "Live Pass / Fail / Blocked tracking per run, module, and environment."),
            ("Defect Management", "Step-level defect creation, Pareto analysis, and CAPA remediation.")
        ]),
        ("ANALYZE & REPORT", TEAL_ACCENT, [
            ("PMO Executive Portal", "SPI, first-pass yield, defect stability, and plan vs. actual timelines."),
            ("QA Analytics Suite", "Cross-milestone trends and predictive quality indicators."),
            ("One-Click Verdicts", "Automated HTML executive email reports + styled Excel/PDF deliverables.")
        ]),
        ("ENTERPRISE PLATFORM", GOLD_ACCENT, [
            ("15 RBAC Roles", "Department-specific access across FA, PM, Dev, QA, and Leadership."),
            ("Full Immutable Audit", "Searchable audit trail tracking every requirement, test, and verdict change."),
            ("On-Prem / Secure Cloud", "Office 365 SMTP, enterprise encryption, and complete data sovereignty.")
        ]),
    ]

    for i, (title, color, features) in enumerate(pillars):
        c_left = Inches(0.8 + i * 2.95)
        add_card(s6, c_left, Inches(2.1), Inches(2.8), Inches(4.7), BG_CARD, BORDER_CARD)

        # Top Accent Header
        top_pill = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left + Inches(0.1), Inches(2.25), Inches(2.6), Inches(0.4))
        top_pill.fill.solid()
        top_pill.fill.fore_color.rgb = BG_CARD_LIGHT
        top_pill.line.color.rgb = color
        top_pill.line.width = Pt(1)
        tf_tp = top_pill.text_frame
        tf_tp.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_tp = tf_tp.paragraphs[0]
        p_tp.text = title
        p_tp.font.size = Pt(11)
        p_tp.font.bold = True
        p_tp.font.color.rgb = color
        p_tp.alignment = PP_ALIGN.CENTER

        # Features List
        c_box = s6.shapes.add_textbox(c_left + Inches(0.12), Inches(2.8), Inches(2.55), Inches(3.9))
        tf_c = c_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = 0

        for f_idx, (fname, fdesc) in enumerate(features):
            p_fn = tf_c.add_paragraph() if f_idx > 0 else tf_c.paragraphs[0]
            p_fn.text = f"✔ {fname}"
            p_fn.font.size = Pt(12)
            p_fn.font.bold = True
            p_fn.font.color.rgb = TEXT_WHITE
            if f_idx > 0: p_fn.space_before = Pt(10)
            
            p_fd = tf_c.add_paragraph()
            p_fd.text = fdesc
            p_fd.font.size = Pt(10)
            p_fd.font.color.rgb = TEXT_MUTED
            p_fd.space_before = Pt(2)

    s6.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / NOTES FOR C-SUITE:\n"
        "\"This slide demonstrates the complete depth of QMPulse across 4 disciplines and 12 core modules. \n"
        "Whether it is milestone planning, test case management, defect analytics, or enterprise compliance, "
        "every team operates within the same governed environment. "
        "Notice the automated Excel generation: it automatically produces Review Effort logs, Pareto analysis, and CAPA corrective action sheets on demand.\""
    )

    # =========================================================================
    # SLIDE 7: ENTERPRISE GOVERNANCE & AUDIT READINESS
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_background(s7)
    add_header(s7, "Governance & Compliance", "Enterprise-Grade Security, RBAC & Audit Readiness", "Built for heavily regulated, security-conscious enterprise IT environments.")

    gov_cards = [
        ("Role-Based Access Control", "15 Granular Roles Across 5 Departments", PURPLE_ACCENT, [
            "Tailored views for PMO, FA, Dev Lead, QA Lead, and Executives.",
            "Strict segregation of duties enforced by code logic.",
            "PMO-only dedicated reporting portal view on login."
        ]),
        ("Immutable Audit Trail", "Complete Operational Traceability", CYAN_ACCENT, [
            "Every requirement change, status flip, and verdict is logged.",
            "Full history timestamping per test step and defect.",
            "Audit-ready exports for regulatory & compliance reviews."
        ]),
        ("Automated PMO Reporting", "One-Click Verdicts & Executive Delivery", TEAL_ACCENT, [
            "HTML email verdicts delivered directly via Office 365 SMTP.",
            "Embedded PDF sign-off certificates with detailed appendices.",
            "Auto-generated Excel workbooks with Pareto & CAPA tabs."
        ]),
        ("Security & Enterprise Hosting", "On-Premises or Private Cloud Ready", EMERALD_ACCENT, [
            "Deployable within your secure corporate network.",
            "PostgreSQL database with JWT authentication (8h expiry).",
            "Zero external data leakage; enterprise API key controls."
        ]),
    ]

    for i, (title, subtitle, color, points) in enumerate(gov_cards):
        row = i // 2
        col = i % 2
        c_left = Inches(0.8 + col * 5.95)
        c_top = Inches(2.1 + row * 2.45)
        add_card(s7, c_left, c_top, Inches(5.75), Inches(2.25), BG_CARD, BORDER_CARD)

        # Card Title
        tbox = s7.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.15), Inches(5.35), Inches(1.95))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.font.name = "Segoe UI"

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.font.name = "Calibri"
        p_sub.space_before = Pt(2)

        for pt in points:
            p_pt = tf.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.size = Pt(10.5)
            p_pt.font.color.rgb = TEXT_MUTED
            p_pt.font.name = "Calibri"
            p_pt.space_before = Pt(3)

    s7.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / NOTES FOR C-SUITE:\n"
        "\"Governance and security are fundamental requirements for our leadership team. \n"
        "- Security & Access: QMPulse enforces 15 granular roles. PMO members get an executive portal; FAs and QA get specialized workspaces.\n"
        "- Compliance: Every test run, status flip, requirement edit, and sign-off is logged in an immutable audit trail.\n"
        "- One-Click Verdicts: Sending release sign-offs to PMO stakeholders takes one click—generating rich HTML emails, PDF certificates, and pre-formatted Excel deliverables with Pareto defect analysis and CAPA tracking.\""
    )

    # =========================================================================
    # SLIDE 8: BUSINESS VALUE & MEASURABLE ROI
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_background(s8)
    add_header(s8, "Business Impact", "Measurable ROI & Delivery Acceleration", "Tangible operational gains delivered directly to the enterprise bottom line.")

    roi_metrics = [
        ("40%", "FASTER TEST DESIGN", "AI Requirement Analyzer and Test Generator cut authoring time by nearly half.", TEAL_ACCENT),
        ("100%", "TRACEABILITY", "Zero orphaned requirements or untracked defects from intake to UAT sign-off.", CYAN_ACCENT),
        ("0", "SPREADSHEET LEAKS", "Eliminates Excel attachment version mismatch and manual consolidation overhead.", INDIGO_ACCENT),
        ("3x", "FASTER REPORTING", "Automated PMO verdict emails and audit Excel generation in under 60 seconds.", EMERALD_ACCENT),
    ]

    for i, (val, title, desc, color) in enumerate(roi_metrics):
        c_left = Inches(0.8 + i * 2.95)
        add_card(s8, c_left, Inches(2.1), Inches(2.8), Inches(2.3), BG_CARD, BORDER_CARD)

        tbox = s8.shapes.add_textbox(c_left + Inches(0.15), Inches(2.25), Inches(2.5), Inches(2.0))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p_v = tf.paragraphs[0]
        p_v.text = val
        p_v.font.size = Pt(28)
        p_v.font.bold = True
        p_v.font.color.rgb = color
        p_v.font.name = "Segoe UI"

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.space_before = Pt(4)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(3)

    # Bottom Qualitative Value Card
    add_card(s8, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.2), BG_CARD_LIGHT, TEAL_ACCENT)
    vbox = s8.shapes.add_textbox(Inches(1.1), Inches(4.85), Inches(11.1), Inches(1.9))
    tf_v = vbox.text_frame
    tf_v.word_wrap = True
    tf_v.margin_left = tf_v.margin_right = tf_v.margin_top = tf_v.margin_bottom = 0

    p_vt = tf_v.paragraphs[0]
    p_vt.text = "STRATEGIC EXECUTIVE ADVANTAGES"
    p_vt.font.size = Pt(12)
    p_vt.font.bold = True
    p_vt.font.color.rgb = TEAL_ACCENT

    val_points = [
        ("Predictable Releases", "Real-time Schedule Performance Index (SPI) and blocker tracking prevent last-minute launch delays."),
        ("Higher First-Pass Yield", "AI validation at the requirement stage prevents defects from cascading into expensive production fixes."),
        ("Operational Agility", "Delivery teams spend time building and verifying value rather than manually filling status reports.")
    ]
    for vt, vd in val_points:
        p_item = tf_v.add_paragraph()
        p_item.text = f"★  {vt}:  {vd}"
        p_item.font.size = Pt(11)
        p_item.font.color.rgb = TEXT_WHITE
        p_item.space_before = Pt(5)

    s8.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / NOTES FOR C-SUITE:\n"
        "\"Let us talk about the business numbers and ROI for executive leadership:\n"
        "1. Efficiency: By generating test cases with Google GenAI, our QA team cuts test design effort by up to 40%.\n"
        "2. Zero Defect Leakage: 100% of requirements map to test executions and Redmine issues.\n"
        "3. Time Savings: PMO reporting cycles that used to take days of manual Excel merging now execute in under 60 seconds with automated verdict emails.\n"
        "4. Predictability: Leadership knows the exact health and stability of every milestone well before go-live.\""
    )

    # =========================================================================
    # SLIDE 9: STRATEGIC ROLLOUT & CALL TO ACTION
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_background(s9)
    add_header(s9, "Implementation & Next Steps", "Strategic Rollout Roadmap & Action Plan", "A phased, low-risk adoption path delivering immediate project value.")

    roadmap_steps = [
        ("Phase 1: Pilot", "Weeks 1 - 2", PURPLE_ACCENT, "Deploy on 2 flagship projects; configure Redmine API sync & role matrices."),
        ("Phase 2: Enablement", "Weeks 3 - 4", INDIGO_ACCENT, "Onboard PMO, FA, Dev, and QA leads; activate AI Requirement Analyzer."),
        ("Phase 3: Scale", "Month 2", CYAN_ACCENT, "Roll out across all enterprise projects; standardize automated PMO verdict delivery."),
        ("Phase 4: Optimization", "Month 3+", EMERALD_ACCENT, "Continuous QA analytics, predictive defect trend modeling & lessons learned review.")
    ]

    for i, (title, timeline, color, desc) in enumerate(roadmap_steps):
        c_left = Inches(0.8 + i * 2.95)
        add_card(s9, c_left, Inches(2.1), Inches(2.8), Inches(2.4), BG_CARD, BORDER_CARD)

        tbox = s9.shapes.add_textbox(c_left + Inches(0.15), Inches(2.25), Inches(2.5), Inches(2.1))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.font.name = "Segoe UI"

        p_time = tf.add_paragraph()
        p_time.text = timeline
        p_time.font.size = Pt(10.5)
        p_time.font.bold = True
        p_time.font.color.rgb = TEXT_WHITE
        p_time.space_before = Pt(3)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(5)

    # Big Executive Call to Action Box
    add_card(s9, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.1), BG_CARD_LIGHT, TEAL_ACCENT)
    cta_box = s9.shapes.add_textbox(Inches(1.2), Inches(4.95), Inches(10.9), Inches(1.8))
    tf_cta = cta_box.text_frame
    tf_cta.word_wrap = True
    tf_cta.margin_left = tf_cta.margin_right = tf_cta.margin_top = tf_cta.margin_bottom = 0

    p_ctat = tf_cta.paragraphs[0]
    p_ctat.text = "THE VERDICT: READY TO TRANSFORM DELIVERY QUALITY"
    p_ctat.font.size = Pt(14)
    p_ctat.font.bold = True
    p_ctat.font.color.rgb = TEAL_ACCENT

    p_ctab = tf_cta.add_paragraph()
    p_ctab.text = "QMPulse is fully built, integrated, and ready to deploy. We invite executive leadership to approve the Phase 1 Pilot rollout starting this week.\n\nThank you — Opening the floor for Questions & Strategic Discussion."
    p_ctab.font.size = Pt(12)
    p_ctab.font.color.rgb = TEXT_WHITE
    p_ctab.space_before = Pt(6)

    s9.notes_slide.notes_text_frame.text = (
        "SPEAKER SCRIPT / NOTES FOR C-SUITE:\n"
        "\"To conclude, our path forward is clear, low-risk, and structured:\n"
        "- We begin with a targeted 2-week pilot on two core projects to validate the workflow with zero operational disruption.\n"
        "- By Month 2, all teams will be unified under one automated quality standard.\n"
        "The software is live, tested, and fully functional today. Thank you for your time, and I now welcome your questions and feedback.\""
    )

    output_path = os.path.join(os.getcwd(), "QMPulse_Executive_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_deck()
